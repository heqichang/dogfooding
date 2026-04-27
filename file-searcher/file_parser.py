import os
from typing import Optional
from constants import TEXT_EXTENSIONS, DOCUMENT_EXTENSIONS
from utils import read_text_file


class FileContentParser:
    @staticmethod
    def parse(file_path: str, max_size: int = 10 * 1024 * 1024) -> Optional[str]:
        if not os.path.exists(file_path) or not os.path.isfile(file_path):
            return None
        
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in TEXT_EXTENSIONS:
            return read_text_file(file_path, max_size)
        elif ext in DOCUMENT_EXTENSIONS:
            return FileContentParser._parse_document(file_path, ext)
        
        return None
    
    @staticmethod
    def _parse_document(file_path: str, ext: str) -> Optional[str]:
        if ext == ".pdf":
            return FileContentParser._parse_pdf(file_path)
        elif ext in (".doc", ".docx"):
            return FileContentParser._parse_word(file_path, ext)
        elif ext in (".xls", ".xlsx"):
            return FileContentParser._parse_excel(file_path, ext)
        elif ext in (".ppt", ".pptx"):
            return FileContentParser._parse_powerpoint(file_path, ext)
        elif ext == ".odt":
            return FileContentParser._parse_odt(file_path)
        elif ext == ".rtf":
            return FileContentParser._parse_rtf(file_path)
        
        return None
    
    @staticmethod
    def _parse_pdf(file_path: str) -> Optional[str]:
        try:
            import pdfplumber
            text_parts = []
            
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
            
            return "\n".join(text_parts) if text_parts else None
        except ImportError:
            return None
        except Exception:
            return None
    
    @staticmethod
    def _parse_word(file_path: str, ext: str) -> Optional[str]:
        if ext == ".docx":
            try:
                from docx import Document
                doc = Document(file_path)
                text_parts = []
                
                for para in doc.paragraphs:
                    text_parts.append(para.text)
                
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            text_parts.append(cell.text)
                
                return "\n".join(text_parts) if text_parts else None
            except ImportError:
                return None
            except Exception:
                return None
        
        elif ext == ".doc":
            try:
                import subprocess
                import platform
                
                if platform.system() == "Windows":
                    try:
                        import win32com.client
                        word = win32com.client.Dispatch("Word.Application")
                        word.Visible = False
                        doc = word.Documents.Open(file_path)
                        text = doc.Content.Text
                        doc.Close(False)
                        word.Quit()
                        return text
                    except ImportError:
                        return None
                    except Exception:
                        try:
                            word.Quit()
                        except:
                            pass
                        return None
                
                return None
            except Exception:
                return None
        
        return None
    
    @staticmethod
    def _parse_excel(file_path: str, ext: str) -> Optional[str]:
        try:
            from openpyxl import load_workbook
            text_parts = []
            
            wb = load_workbook(file_path, read_only=True, data_only=True)
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")
                
                for row in ws.iter_rows(values_only=True):
                    row_text = " ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            wb.close()
            return "\n".join(text_parts) if text_parts else None
        except ImportError:
            return None
        except Exception:
            return None
    
    @staticmethod
    def _parse_powerpoint(file_path: str, ext: str) -> Optional[str]:
        if ext == ".pptx":
            try:
                from pptx import Presentation
                text_parts = []
                
                prs = Presentation(file_path)
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    text_parts.append(f"Slide {slide_num}:")
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                        
                        if hasattr(shape, "table"):
                            for row in shape.table.rows:
                                row_text = " ".join(cell.text for cell in row.cells if cell.text.strip())
                                if row_text:
                                    text_parts.append(row_text)
                
                return "\n".join(text_parts) if text_parts else None
            except ImportError:
                return None
            except Exception:
                return None
        
        return None
    
    @staticmethod
    def _parse_odt(file_path: str) -> Optional[str]:
        try:
            import zipfile
            from xml.etree import ElementTree as ET
            
            with zipfile.ZipFile(file_path, 'r') as zf:
                if 'content.xml' not in zf.namelist():
                    return None
                
                with zf.open('content.xml') as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    
                    namespaces = {
                        'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0',
                        'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0',
                    }
                    
                    text_parts = []
                    
                    for elem in root.iter():
                        if elem.tag.endswith('}p') or elem.tag.endswith('}h'):
                            if elem.text:
                                text_parts.append(elem.text)
                            for child in elem.iter():
                                if child.tail:
                                    text_parts.append(child.tail)
                    
                    return '\n'.join(text_parts) if text_parts else None
        except ImportError:
            return None
        except Exception:
            return None
    
    @staticmethod
    def _parse_rtf(file_path: str) -> Optional[str]:
        try:
            import subprocess
            import platform
            
            if platform.system() == "Windows":
                try:
                    import win32com.client
                    word = win32com.client.Dispatch("Word.Application")
                    word.Visible = False
                    doc = word.Documents.Open(file_path)
                    text = doc.Content.Text
                    doc.Close(False)
                    word.Quit()
                    return text
                except ImportError:
                    pass
                except Exception:
                    try:
                        word.Quit()
                    except:
                        pass
            
            return None
        except Exception:
            return None
